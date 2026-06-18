"""
Build an editable PPTX from the HTML slide design.
All text, shapes, and images are native editable elements.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).parent

# ===== Colors (matching light theme) =====
C_TEXT       = RGBColor(0x1a, 0x25, 0x40)  # dark navy text
C_TEXT_SUB   = RGBColor(0x5a, 0x6a, 0x8c)  # secondary text
C_TEXT_MUTED = RGBColor(0x8a, 0x9a, 0xb2)  # muted
C_CYAN       = RGBColor(0x00, 0x99, 0xcc)  # primary cyan
C_CYAN_DARK  = RGBColor(0x00, 0x78, 0xa8)  # code text
C_PURPLE     = RGBColor(0x6c, 0x2b, 0xff)  # data sovereignty
C_PURPLE_D   = RGBColor(0x5a, 0x1f, 0xd8)
C_BORDER     = RGBColor(0xc0, 0xcb, 0xde)  # neutral border
C_BORDER_CYAN= RGBColor(0x66, 0xc4, 0xe0)
C_BORDER_PUR = RGBColor(0xae, 0x95, 0xf0)
C_BG_PALE    = RGBColor(0xf5, 0xf8, 0xfc)
C_BG_CYAN    = RGBColor(0xe6, 0xf5, 0xfa)
C_BG_PURPLE  = RGBColor(0xf0, 0xeb, 0xff)
C_WHITE      = RGBColor(0xff, 0xff, 0xff)

# ===== Slide setup =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Virtual 1600x900 canvas → EMU
def emu_x(px): return int(px * SLIDE_W / 1600)
def emu_y(px): return int(px * SLIDE_H / 900)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# Background (white)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = C_WHITE


def add_text(x, y, w, h, text, font_size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="PingFang SC"):
    tb = slide.shapes.add_textbox(emu_x(x), emu_y(y), emu_x(w), emu_y(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_runs(x, y, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="PingFang SC"):
    """parts: list of (text, font_size, bold, color)"""
    tb = slide.shapes.add_textbox(emu_x(x), emu_y(y), emu_x(w), emu_y(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, sz, bold, color) in enumerate(parts):
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(x, y, w, h, fill=C_WHITE, line=C_BORDER, line_w=0.75,
             rounded=False, corner=0.06):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, emu_x(x), emu_y(y), emu_x(w), emu_y(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if rounded:
        # adjust corner radius
        s.adjustments[0] = corner
    s.shadow.inherit = False
    return s


def add_image(x, y, w, h, path):
    return slide.shapes.add_picture(str(path), emu_x(x), emu_y(y), emu_x(w), emu_y(h))


# ============= LAYOUT =============
# Slide: 1600 x 900
# Padding: 44 left/right, 28 top, 32 bottom

# --- Title ---
add_runs(44, 28, 1512, 50, [
    ("AI 工具链可观测 ", 26, True, C_TEXT),
    ("— 看得见 · 可追溯 · 数据自主", 26, True, C_CYAN),
], anchor=MSO_ANCHOR.TOP)

# --- Subtitle ---
add_runs(44, 78, 1512, 28, [
    ("使用", 12, False, C_TEXT_SUB),
    ("阿里云 AI 工具链", 12, True, C_CYAN_DARK),
    (" · 让 AI 替开发者 / 运维操作阿里云,从“我相信它”升级到 ", 12, False, C_TEXT_SUB),
    ("看得见", 12, True, C_CYAN_DARK),
    ("  ", 12, False, C_TEXT_SUB),
    ("可追溯", 12, True, C_CYAN_DARK),
    ("  ", 12, False, C_TEXT_SUB),
    ("可量化", 12, True, C_CYAN_DARK),
    ("  ", 12, False, C_TEXT_SUB),
    ("可审计", 12, True, C_CYAN_DARK),
    ("  ", 12, False, C_TEXT_SUB),
    ("能合规", 12, True, C_CYAN_DARK),
])

# --- Body area ---
# Body starts at y=130, height=740
# Left arch: x=44, y=130, w=1124, h=740
# Right values: x=1196, y=130, w=360, h=740

# === LEFT: Arch container ===
add_rect(44, 130, 1124, 740, fill=RGBColor(0xfb, 0xfc, 0xfe), line=C_BORDER, line_w=0.5, rounded=True, corner=0.015)

# --- Operation flow (4 nodes + arrows) ---
# inside arch: x_start=60, y_start=144
# Nodes total width: 1092 - arrow widths
# User: 134 fixed
# Arrows: ~14px each, 3 arrows
# Available: 1092 - 134 - 14*3 - 7*3(gaps) = 895
# Agent:Tools:Assets = 1.05:1.4:1.2 = 3.65
NODE_Y = 148
NODE_H = 92
USER_W = 134
ARROW_W = 16
GAP = 8
avail = 1092 - USER_W - ARROW_W*3 - GAP*3
AGENT_W = int(avail * 1.05/3.65)
TOOLS_W = int(avail * 1.4/3.65)
ASSETS_W = avail - AGENT_W - TOOLS_W

# User node (centered text)
ux = 60
add_rect(ux, NODE_Y, USER_W, NODE_H, fill=C_WHITE, line=C_BORDER, rounded=True, corner=0.08)
# Icon placeholder (circle representing person)
icon_d = 24
icon_x = ux + (USER_W - icon_d) // 2
icon_y = NODE_Y + 12
icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu_x(icon_x), emu_y(icon_y), emu_x(icon_d), emu_y(icon_d))
icon.fill.background()
icon.line.color.rgb = C_TEXT_SUB
icon.line.width = Pt(1.2)
add_text(ux+4, NODE_Y+40, USER_W-8, 18, "用户 / 开发者 / 运维", font_size=10, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
add_text(ux+4, NODE_Y+58, USER_W-8, 14, "用云需求", font_size=8.5, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)

# Arrow
ax = ux + USER_W + GAP
arrow_y = NODE_Y + (NODE_H // 2) - 5
add_text(ax, arrow_y, ARROW_W, 14, "▶", font_size=10, bold=True, color=C_BORDER, align=PP_ALIGN.CENTER)

# Agent node
agx = ax + ARROW_W + GAP
add_rect(agx, NODE_Y, AGENT_W, NODE_H, fill=RGBColor(0xf7, 0xf3, 0xff), line=C_BORDER_PUR, line_w=0.8, rounded=True, corner=0.08)
# Agent header
add_text(agx+12, NODE_Y+10, AGENT_W-20, 18, "■  Agent", font_size=11, bold=True, color=C_TEXT)
# Override the "■" with purple color by re-creating
agent_head_tb = slide.shapes[-1]
agent_head_tb.text_frame.paragraphs[0].runs[0].font.color.rgb = C_PURPLE
# Actually since we can't easily mix colors in one run, let's use add_runs
# Remove and redo
sp = agent_head_tb._element
sp.getparent().remove(sp)
add_runs(agx+12, NODE_Y+10, AGENT_W-20, 18, [
    ("■  ", 11, True, C_PURPLE),
    ("Agent", 11, True, C_TEXT),
])
# Client chips (text only, separated by " · ")
add_text(agx+12, NODE_Y+34, AGENT_W-20, 50,
         "Claude Code  ·  Codex  ·  Qoder Work  ·  …",
         font_size=9, bold=True, color=C_PURPLE_D)

# Arrow 2
ax2 = agx + AGENT_W + GAP
add_text(ax2, arrow_y, ARROW_W, 14, "▶", font_size=10, bold=True, color=C_BORDER, align=PP_ALIGN.CENTER)

# Tools node
tx = ax2 + ARROW_W + GAP
add_rect(tx, NODE_Y, TOOLS_W, NODE_H, fill=RGBColor(0xe9, 0xf6, 0xfb), line=C_BORDER_CYAN, line_w=0.8, rounded=True, corner=0.08)
add_runs(tx+12, NODE_Y+10, TOOLS_W-20, 18, [
    ("◆  ", 11, True, C_CYAN),
    ("阿里云 AI 工具链", 11, True, C_TEXT),
])
add_text(tx+12, NODE_Y+34, TOOLS_W-20, 50,
         "CLI  ·  SDK  ·  MCP  ·  Terraform  ·  Skill  ·  Plugin",
         font_size=9, bold=True, color=C_CYAN_DARK)

# Arrow 3
ax3 = tx + TOOLS_W + GAP
add_text(ax3, arrow_y, ARROW_W, 14, "▶", font_size=10, bold=True, color=C_BORDER, align=PP_ALIGN.CENTER)

# Assets node
asx = ax3 + ARROW_W + GAP
add_rect(asx, NODE_Y, ASSETS_W, NODE_H, fill=C_WHITE, line=C_BORDER, line_w=0.8, rounded=True, corner=0.08)
add_runs(asx+12, NODE_Y+10, ASSETS_W-20, 18, [
    ("☁  ", 11, True, C_TEXT_SUB),
    ("阿里云资产", 11, True, C_TEXT),
])
add_text(asx+12, NODE_Y+34, ASSETS_W-20, 50,
         "ECS · RDS · VPC · OSS · ACK · RAM",
         font_size=9, bold=True, color=C_TEXT_SUB)


# --- trace-line ---
trace_y = NODE_Y + NODE_H + 6
add_text(60, trace_y, 1092, 18,
         "▼   t r a c i n g   全 程 旁 路 记 录   ·   不 打 扰   A g e n t   工 作   ▼",
         font_size=8, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)

# --- panels-head ---
ph_y = trace_y + 22
add_runs(60, ph_y, 1092, 18, [
    ("同一份 Trace · ", 9, True, C_TEXT_SUB),
    ("多客户端开箱即用", 9, True, C_CYAN),
])

# --- 4 panels (2x2 grid) ---
PANELS_Y = ph_y + 24
PANELS_X = 60
PANELS_W = 1092
PANELS_H = 712 - (PANELS_Y - 144)  # remaining vertical
P_GAP = 8

CELL_W = (PANELS_W - P_GAP) // 2
CELL_H = (PANELS_H - P_GAP) // 2

PANEL_HEAD_H = 26

panels_data = [
    ("C", "Claude Code", "本地 · 默认开启", "本地可观测-claudecode客户端.png", "cyan"),
    ("X", "Codex",        "本地 · 默认开启", "本地可观测-codex客户端.png",      "cyan"),
    ("Q", "Qoder Work",   "本地 · 默认开启", "本地可观测-qoderwork客户端.png",  "cyan"),
    ("☁","客户自有数仓", "远程 · 可选 BYOC", "远程可观测上报.png",              "purple"),
]
badge_colors = {
    "cyan":   (C_BG_CYAN, C_BORDER_CYAN, C_CYAN_DARK),
    "purple": (C_BG_PURPLE, C_BORDER_PUR, C_PURPLE_D),
}
em_badge_colors = {
    "C":  RGBColor(0xff, 0x7a, 0x3d),
    "X":  RGBColor(0x2c, 0x2c, 0x2c),
    "Q":  RGBColor(0x4a, 0x8f, 0xff),
    "☁": C_PURPLE,
}

for i, (em, name, badge_text, img, color_key) in enumerate(panels_data):
    col = i % 2
    row = i // 2
    px = PANELS_X + col * (CELL_W + P_GAP)
    py = PANELS_Y + row * (CELL_H + P_GAP)

    badge_bg, panel_border, badge_fg = badge_colors[color_key]
    # Panel container
    add_rect(px, py, CELL_W, CELL_H, fill=C_WHITE, line=panel_border, line_w=0.8, rounded=True, corner=0.02)
    # Panel head bar
    add_rect(px+1, py+1, CELL_W-2, PANEL_HEAD_H, fill=C_BG_PALE, line=None)
    # Em badge (small colored box w/ letter)
    em_box_d = 16
    em_x = px + 8
    em_y = py + 1 + (PANEL_HEAD_H - em_box_d) // 2
    em_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      emu_x(em_x), emu_y(em_y),
                                      emu_x(em_box_d), emu_y(em_box_d))
    em_rect.fill.solid()
    em_rect.fill.fore_color.rgb = em_badge_colors[em]
    em_rect.line.fill.background()
    em_rect.adjustments[0] = 0.15
    em_rect.shadow.inherit = False
    # Letter inside badge
    add_text(em_x, em_y, em_box_d, em_box_d, em, font_size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Panel name
    add_text(em_x + em_box_d + 6, py + 4, 200, PANEL_HEAD_H - 6, name, font_size=10, bold=True, color=C_TEXT)
    # Badge on right
    badge_w = 92
    badge_h = 14
    badge_x = px + CELL_W - badge_w - 8
    badge_y = py + 1 + (PANEL_HEAD_H - badge_h) // 2
    add_rect(badge_x, badge_y, badge_w, badge_h, fill=badge_bg, line=panel_border, line_w=0.5, rounded=True, corner=0.4)
    add_text(badge_x, badge_y, badge_w, badge_h, badge_text, font_size=7.5, bold=True, color=badge_fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Screenshot
    img_path = HERE / img
    img_y = py + PANEL_HEAD_H + 4
    img_h = CELL_H - PANEL_HEAD_H - 8
    img_w = CELL_W - 8
    img_x = px + 4
    # We use add_picture with full cell area (PPT will fit)
    add_image(img_x, img_y, img_w, img_h, img_path)


# === RIGHT: 5 Value cards ===
VAL_X = 1196
VAL_Y = 130
VAL_W = 360
VAL_H = 740
CARD_GAP = 10
CARD_H = (VAL_H - 4 * CARD_GAP) // 5

cards = [
    ("看得见", None, "每个 turn 一棵 span 树 — LLM 思考 → 工具选择 → 阿里云 OpenAPI 全链路", "cyan"),
    ("可追溯", None, "出错时展开错误堆栈,RequestId 定位阿里云控制台 / ActionTrail", "cyan"),
    ("可量化", None, "每会话 turn / tool / skill / 成功率 / token,本地直查 P95 与错误率", "cyan"),
    ("可审计", None, "操作五元组留痕:RAM 身份 · 时间戳 · OpenAPI Action · 资源 ID · RequestId", "cyan"),
    ("能合规 · 数据自主", "DATA SOVEREIGNTY",
     "AK / SK / PII 自动脱敏 · 可上报至 SLS / ARMS / 客户自建数仓", "purple"),
]

for idx, (title, en, desc, color_key) in enumerate(cards):
    cy = VAL_Y + idx * (CARD_H + CARD_GAP)
    if color_key == "purple":
        accent = C_PURPLE
        icon_bg = C_BG_PURPLE
        icon_border = C_BORDER_PUR
    else:
        accent = C_CYAN
        icon_bg = C_BG_CYAN
        icon_border = C_BORDER_CYAN
    # Card container
    add_rect(VAL_X, cy, VAL_W, CARD_H, fill=C_WHITE, line=C_BORDER, line_w=0.5, rounded=True, corner=0.05)
    # Left accent bar (thin rect)
    add_rect(VAL_X, cy, 3, CARD_H, fill=accent, line=None)
    # Icon box (left)
    icon_size = 38
    icon_x = VAL_X + 14
    icon_y = cy + (CARD_H - icon_size) // 2
    add_rect(icon_x, icon_y, icon_size, icon_size, fill=icon_bg, line=icon_border, line_w=0.6, rounded=True, corner=0.15)
    # Icon as a small symbol inside (textual placeholder)
    glyphs = ["◉", "○", "▦", "▤", "◈"]
    add_text(icon_x, icon_y, icon_size, icon_size, glyphs[idx], font_size=16, bold=True, color=accent, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name="Helvetica")
    # Title
    text_x = icon_x + icon_size + 12
    text_w = VAL_X + VAL_W - text_x - 12
    title_y = cy + 12 if en is None else cy + 8
    add_text(text_x, title_y, text_w, 22, title, font_size=12.5, bold=True, color=C_TEXT)
    # English (only card 5)
    if en:
        add_text(text_x, cy + 28, text_w, 14, en, font_size=8, bold=True, color=accent)
        desc_y = cy + 46
    else:
        desc_y = cy + 36
    # Description
    add_text(text_x, desc_y, text_w, CARD_H - (desc_y - cy) - 8, desc,
             font_size=9.5, color=C_TEXT_SUB)


# === Footer ===
add_text(44, 870, 400, 20, "阿里云 Agent Toolkit", font_size=8.5, color=C_TEXT_MUTED)
add_runs(900, 870, 656, 20, [
    ("多客户端开箱即用     本地默认归档     远程上报至 ", 8.5, False, C_TEXT_MUTED),
    ("客户自有数仓", 8.5, True, C_CYAN),
], align=PP_ALIGN.RIGHT)


# === Save ===
out = HERE / "observability-slide.pptx"
prs.save(out)
print(f"Saved: {out} ({out.stat().st_size // 1024} KB)")
